"""
MCQ Generator - Generate multiple choice questions from documents using Gemini API.
"""

import pathlib
import sqlite3
import sys
from typing import List

from dotenv import load_dotenv
from google import genai
from google.genai import types
from pptx import Presentation
from pydantic import BaseModel, Field

load_dotenv()
# ============================================================================
# Pydantic Models for Structured Output
# ============================================================================


class Option(BaseModel):
    """A single option for a multiple choice question."""

    text: str = Field(description="The text content of this option.")
    is_correct: bool = Field(
        description="True if this is the correct answer, False otherwise."
    )


class Question(BaseModel):
    """A single multiple choice question with its options."""

    text: str = Field(description="The question text.")
    options: List[Option] = Field(
        description="List of 4 answer options. Exactly one must be correct.",
        min_length=4,
        max_length=4,
    )


class MCQResponse(BaseModel):
    """Response containing all generated multiple choice questions."""

    questions: List[Question] = Field(
        description="List of multiple choice questions generated from the document."
    )


# ============================================================================
# Database Operations
# ============================================================================


def init_database(db_path: str = "mcqs.db") -> sqlite3.Connection:
    """Initialize the SQLite database with required tables."""
    conn = sqlite3.connect(db_path)
    cursor = conn.cursor()

    # Create questions table
    cursor.execute(
        """
        CREATE TABLE IF NOT EXISTS questions (
            question_id INTEGER PRIMARY KEY AUTOINCREMENT,
            text TEXT NOT NULL
        )
    """
    )

    # Create options table
    cursor.execute(
        """
        CREATE TABLE IF NOT EXISTS options (
            option_id INTEGER PRIMARY KEY AUTOINCREMENT,
            question_id INTEGER NOT NULL,
            text TEXT NOT NULL,
            is_correct INTEGER NOT NULL,
            FOREIGN KEY (question_id) REFERENCES questions(question_id)
        )
    """
    )

    conn.commit()
    return conn


def clear_database(conn: sqlite3.Connection) -> None:
    """Clear all existing data from the database."""
    cursor = conn.cursor()
    cursor.execute("DELETE FROM options")
    cursor.execute("DELETE FROM questions")
    cursor.execute(
        "DELETE FROM sqlite_sequence WHERE name='questions' OR name='options'"
    )
    conn.commit()


def save_questions_to_db(conn: sqlite3.Connection, mcq_response: MCQResponse) -> None:
    """Save generated questions and options to the database."""
    cursor = conn.cursor()

    for question in mcq_response.questions:
        # Insert question
        cursor.execute("INSERT INTO questions (text) VALUES (?)", (question.text,))
        question_id = cursor.lastrowid

        # Insert options
        for option in question.options:
            cursor.execute(
                "INSERT INTO options (question_id, text, is_correct) VALUES (?, ?, ?)",
                (question_id, option.text, 1 if option.is_correct else 0),
            )

    conn.commit()


# ============================================================================
# Gemini API Integration
# ============================================================================

MCQ_SYSTEM_PROMPT = """You are an expert educational assessment creator specializing in generating high-quality multiple choice questions (MCQs).

Your task is to analyze the provided document and create comprehensive multiple choice questions that test understanding of the key concepts, facts, and ideas presented.

Guidelines for creating MCQs:
1. Each question must have exactly 4 options (A, B, C, D)
2. Exactly one option must be correct for each question
3. All incorrect options (distractors) should be plausible but clearly wrong
4. Questions should cover different aspects and difficulty levels of the material
5. Questions should be clear, unambiguous, and test meaningful understanding
6. Avoid trivial questions or those that can be answered without reading the document
7. Do NOT include explanations - only the question and options
8. Make distractors challenging but fair - they should require actual knowledge to distinguish from the correct answer
9. Vary question types: recall, comprehension, application, and analysis
10. Ensure options are similar in length and structure to avoid giving away the answer"""


def get_mime_type(file_path: pathlib.Path) -> str:
    """Determine the MIME type based on file extension."""
    extension = file_path.suffix.lower()
    mime_types = {
        ".pdf": "application/pdf",
        ".txt": "text/plain",
        ".md": "text/markdown",
        ".html": "text/html",
        ".htm": "text/html",
        ".json": "application/json",
        ".xml": "application/xml",
        ".csv": "text/csv",
        ".py": "text/x-python",
        ".js": "text/javascript",
        ".ts": "text/typescript",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    }
    return mime_types.get(extension, "text/plain")


def extract_text_from_pptx(file_path: pathlib.Path) -> str:
    """Extract all text content from a PowerPoint file."""
    prs = Presentation(str(file_path))
    text_parts = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = [f"--- Slide {slide_num} ---"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in paragraph.runs).strip()
                    if text:
                        slide_texts.append(text)
            # Also extract text from tables
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            slide_texts.append(cell.text.strip())
        if len(slide_texts) > 1:  # More than just the slide header
            text_parts.append("\n".join(slide_texts))

    return "\n\n".join(text_parts)


def prepare_file_content(file_path: str) -> types.Part:
    """
    Prepare file content as a Part for the Gemini API.

    Args:
        file_path: Path to the document file

    Returns:
        types.Part containing the file content
    """
    filepath = pathlib.Path(file_path)
    if not filepath.exists():
        raise FileNotFoundError(f"File not found: {file_path}")

    mime_type = get_mime_type(filepath)

    # Handle PPTX files by extracting text (Gemini doesn't support PPTX natively)
    if filepath.suffix.lower() == ".pptx":
        extracted_text = extract_text_from_pptx(filepath)
        file_content = extracted_text.encode("utf-8")
        mime_type = "text/plain"
    else:
        file_content = filepath.read_bytes()

    return types.Part.from_bytes(data=file_content, mime_type=mime_type)


def generate_mcqs_from_files(
    file_paths: List[str], num_questions: int = 10
) -> MCQResponse:
    """
    Generate MCQs from one or more document files using Gemini API.

    Args:
        file_paths: List of paths to document files
        num_questions: Number of questions to generate

    Returns:
        MCQResponse containing generated questions
    """
    if not file_paths:
        raise ValueError("At least one file path must be provided")

    # Initialize Gemini client
    client = genai.Client()

    # Prepare all file contents
    file_parts = [prepare_file_content(fp) for fp in file_paths]

    # Create the prompt
    doc_word = "document" if len(file_paths) == 1 else "documents"
    prompt = f"""Analyze the provided {doc_word} and generate exactly {num_questions} multiple choice questions.

Each question should:
- Test a meaningful concept from the {doc_word}
- Have exactly 4 options with only one correct answer
- Be clear and unambiguous
- Not include any explanation

Generate diverse questions covering different topics and difficulty levels from the {doc_word}."""

    # Build contents: all files followed by the prompt
    contents = file_parts + [prompt]

    # Make API call with structured output
    response = client.models.generate_content(
        model="gemini-2.5-pro",
        contents=contents,
        config=types.GenerateContentConfig(
            system_instruction=MCQ_SYSTEM_PROMPT,
            response_mime_type="application/json",
            response_json_schema=MCQResponse.model_json_schema(),
            thinking_config=types.ThinkingConfig(
                thinking_budget=32000,
            ),
        ),
    )

    # Parse and return the response
    if not response.text:
        raise ValueError("No response received from Gemini API.")

    mcq_response = MCQResponse.model_validate_json(response.text)
    return mcq_response


# ============================================================================
# Main Entry Point
# ============================================================================


def main():
    """Main function to generate MCQs from one or more files."""
    if len(sys.argv) < 2:
        print("Usage: python main.py <file_path> [file_path2 ...] [-n num_questions]")
        print("Examples:")
        print("  python main.py document.pdf")
        print("  python main.py document.pdf -n 15")
        print("  python main.py doc1.pdf doc2.pdf slides.pptx -n 20")
        sys.exit(1)

    # Parse arguments: files and optional -n flag for num_questions
    args = sys.argv[1:]
    num_questions = 10
    file_paths = []

    i = 0
    while i < len(args):
        if args[i] == "-n" and i + 1 < len(args):
            num_questions = int(args[i + 1])
            i += 2
        else:
            file_paths.append(args[i])
            i += 1

    if not file_paths:
        print("Error: At least one file path is required.")
        sys.exit(1)

    file_word = "file" if len(file_paths) == 1 else "files"
    print(f"Generating {num_questions} MCQs from {len(file_paths)} {file_word}:")
    for fp in file_paths:
        print(f"  - {fp}")

    # Initialize database
    conn = init_database()
    clear_database(conn)

    # Generate MCQs
    mcq_response = generate_mcqs_from_files(file_paths, num_questions)

    # Save to database
    save_questions_to_db(conn, mcq_response)

    print(f"Successfully generated {len(mcq_response.questions)} questions!")
    print("Questions saved to mcqs.db")

    conn.close()


if __name__ == "__main__":
    main()
